import os
import streamlit as st
import pandas as pd
from datasets import load_dataset

from sentence_transformers import SentenceTransformer, util
import docx2txt
from pptx import Presentation
import fitz  # PyMuPDF

from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, pipeline

from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.llms.huggingface_pipeline import HuggingFacePipeline
from langchain_core.prompts import PromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser

# --------------------------------------------------
# BASE PATH (works in Colab, local, Streamlit Cloud)
# --------------------------------------------------
BASE_PATH = os.path.dirname(os.path.abspath(__file__))

# --------------------------------------------------
# GROQ API KEY
# --------------------------------------------------
def get_groq_api_key():
    if "GROQ_API_KEY" in os.environ:
        return os.environ["GROQ_API_KEY"]
    return st.secrets["GROQ_API_KEY"]

GROQ_API_KEY = get_groq_api_key()

# --------------------------------------------------
# STREAMLIT CONFIG
# --------------------------------------------------
st.set_page_config(
    page_title="AI Career Assistant",
    layout="wide",
    page_icon="💼"
)

# --------------------------------------------------
# LOAD DATASET
# --------------------------------------------------
@st.cache_data(show_spinner=False)
def load_data():
    dataset = load_dataset("DevilsLord/It_job_roles_skills_certifications")
    df = dataset["train"].to_pandas()
    df["combined"] = (
        df["Job Description"].fillna("")
        + " "
        + df["Skills"].fillna("")
        + " "
        + df["Certifications"].fillna("")
    )
    return df

# --------------------------------------------------
# LOAD RESUME MODEL
# --------------------------------------------------
@st.cache_resource(show_spinner=False)
def load_resume_model():
    return SentenceTransformer("all-MiniLM-L6-v2", device="cpu")

# --------------------------------------------------
# LOAD EMBEDDINGS + FAISS (SAFE)
# --------------------------------------------------
@st.cache_resource(show_spinner=False)
def load_embeddings_and_faiss():
    emb_model = HuggingFaceEmbeddings(
        model_name="sentence-transformers/paraphrase-MiniLM-L3-v2"
    )

    faiss_path = os.path.join(BASE_PATH, "faiss_index")

    # ✅ LOAD EXISTING FAISS (SAFE DESERIALIZATION)
    if os.path.exists(faiss_path):
        vect = FAISS.load_local(
            faiss_path,
            embeddings=emb_model,
            allow_dangerous_deserialization=True
        )
        return emb_model, vect

    # 🔨 BUILD FAISS IF NOT PRESENT
    df = load_data()
    vect = FAISS.from_texts(df["combined"].tolist(), embedding=emb_model)
    vect.save_local(faiss_path)

    return emb_model, vect

# --------------------------------------------------
# BUILD RAG CHAIN
# --------------------------------------------------
@st.cache_resource(show_spinner=False)
def build_rag_chain():
    from langchain_groq import ChatGroq

    _, vectorstore = load_embeddings_and_faiss()

    llm = ChatGroq(
        groq_api_key=GROQ_API_KEY,
        model_name="llama-3.1-8b-instant",
        temperature=0.2
    )

    prompt = PromptTemplate.from_template(
        """
You are an expert IT career assistant.
Answer clearly and professionally using ONLY the context.
If the answer is not present, say "I don't have that information."

Context:
{context}

Question:
{question}

Answer:
"""
    )

    retriever = vectorstore.as_retriever(search_kwargs={"k": 4})

    chain = (
        {"context": retriever, "question": RunnablePassthrough()}
        | prompt
        | llm
        | StrOutputParser()
    )

    return chain

# --------------------------------------------------
# TEXT EXTRACTION
# --------------------------------------------------
def extract_text(file):
    ext = file.name.split(".")[-1].lower()

    if ext == "txt":
        return file.read().decode("utf-8", errors="ignore")

    if ext == "pdf":
        doc = fitz.open(stream=file.read(), filetype="pdf")
        return "\n".join(page.get_text() for page in doc)

    if ext == "docx":
        return docx2txt.process(file)

    if ext == "pptx":
        text = ""
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    return ""

# --------------------------------------------------
# INITIALIZE MODELS
# --------------------------------------------------
with st.spinner("Initializing AI models..."):
    df = load_data()
    resume_model = load_resume_model()
    rag_chain = build_rag_chain()

# --------------------------------------------------
# UI
# --------------------------------------------------
st.sidebar.title("🌙 Navigation")
page = st.sidebar.radio(
    "Go to",
    ["🏠 Home", "📄 Resume Matcher", "💬 Chatbot"],
    label_visibility="collapsed"
)

# ---------------- HOME ----------------
if page == "🏠 Home":
    st.title("🧠 AI-Powered Career Assistant")
    st.markdown(
        """
        This application helps you:
        - Match resumes with IT job roles
        - Chat with an AI assistant about IT careers, skills, and certifications
        """
    )

# ---------------- RESUME MATCHER ----------------
elif page == "📄 Resume Matcher":
    st.title("📄 AI Resume Matcher")

    uploaded = st.file_uploader(
        "Upload Resume (PDF, DOCX, PPTX, TXT)",
        type=["pdf", "docx", "pptx", "txt"]
    )

    if uploaded:
        resume_text = extract_text(uploaded).replace("\n", " ")
        resume_emb = resume_model.encode(resume_text, convert_to_tensor=True)
        job_embs = resume_model.encode(df["combined"].tolist(), convert_to_tensor=True)

        sims = util.cos_sim(resume_emb, job_embs)[0].cpu().numpy()
        df["Similarity"] = sims

        top = df.sort_values("Similarity", ascending=False).head(10)

        for _, row in top.iterrows():
            st.markdown(
                f"**{row['Job Title']}** — Match Score: `{row['Similarity']:.2f}`"
            )

# ---------------- CHATBOT ----------------
elif page == "💬 Chatbot":
    st.title("💬 IT Career Chatbot")

    if "messages" not in st.session_state:
        st.session_state.messages = [
            {
                "role": "assistant",
                "content": "Hi! Ask me about IT careers, skills, or certifications."
            }
        ]

    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.write(msg["content"])

    if user_input := st.chat_input("Ask a question"):
        st.session_state.messages.append(
            {"role": "user", "content": user_input}
        )

        with st.chat_message("assistant"):
            response = rag_chain.invoke(user_input)
            st.write(response)

        st.session_state.messages.append(
            {"role": "assistant", "content": response}
        )
